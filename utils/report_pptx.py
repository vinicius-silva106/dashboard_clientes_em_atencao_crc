import io
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def generate_pptx_report(kpis: list, esc_df: pd.DataFrame, trib_df: pd.DataFrame) -> bytes:
    """Gera um arquivo PPTX em memória com os dados executivos para importação no Google Slides."""
    prs = Presentation()
    
    # CORES TEMA
    COLOR_TEAL = RGBColor(5, 150, 105)
    COLOR_NAVY = RGBColor(27, 38, 59)
    COLOR_NEUTRAL = RGBColor(100, 116, 139)

    # ---------------------------------------------------------
    # SLIDE 1: CAPA
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[0]) # Title slide
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Dashboard Executivo CRC"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_NAVY
    
    subtitle.text = "Panorama de Clientes em Atenção\nRelatório Gerado Automaticamente"
    subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_TEAL

    # ---------------------------------------------------------
    # SLIDE 2: KPIs GERAIS
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
    slide.shapes.title.text = "Resumo Executivo (KPIs)"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = COLOR_NAVY
    
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    for kpi in kpis:
        p = tf.add_paragraph()
        lbl = kpi.get("label", "Dado")
        cur = kpi.get("current", 0)
        prev = kpi.get("prev", 0)
        
        p.text = f"• {lbl}: {cur} (Semana Anterior: {prev})"
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_NEUTRAL

    # ---------------------------------------------------------
    # SLIDE 3: ESCALATIONS
    # ---------------------------------------------------------
    if not esc_df.empty:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Lista de Casos - Escalations"
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = COLOR_NAVY
        
        limit = min(8, len(esc_df))
        rows, cols = limit + 1, 3
        left, top, width, height = Inches(0.5), Inches(1.8), Inches(9), Inches(0.6)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        headers = ["Cliente", "Produto", "HealthScore"]
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = COLOR_TEAL
            
        for idx in range(limit):
            row_data = esc_df.iloc[idx]
            table.cell(idx+1, 0).text = str(row_data.get("Cliente", ""))
            table.cell(idx+1, 1).text = str(row_data.get("Produto", ""))
            
            hs = row_data.get("HealthScore")
            table.cell(idx+1, 2).text = f"{hs:.1f}%" if pd.notnull(hs) else "N/D"

    # ---------------------------------------------------------
    # COMPILAR PARA BYTES
    # ---------------------------------------------------------
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io.read()
